import os
import sys
import argparse
import pandas as pd
import pdfplumber
from docx import Document
from pptx import Presentation
from typing import List, Optional

class ParserError(Exception):
    """Custom exception for parsing errors."""
    pass

def parse_excel(file_path: str) -> str:
    """
    Reads an Excel file and converts all its sheets to Markdown tables.
    """
    try:
        # Read all sheets. sheet_name=None returns a dictionary of dataframes.
        sheets = pd.read_excel(file_path, sheet_name=None)
        all_markdown = []
        for sheet_name, df in sheets.items():
            # Convert to markdown table using to_markdown() which uses tabulate under the hood.
            all_markdown.append(f"## Sheet: {sheet_name}\n\n{df.to_markdown(index=False)}")
        return "\n\n".join(all_markdown)
    except ImportError:
        raise ParserError("The 'tabulate' library is required for Excel to Markdown conversion. Please install it via 'pip install tabulate'.")
    except (ValueError, FileNotFoundError, PermissionError) as e:
        raise ParserError(f"File error parsing Excel file {file_path}: {str(e)}")
    except Exception as e:
        raise ParserError(f"Unexpected error parsing Excel file {file_path}: {str(e)}")

def parse_text(file_path: str) -> str:
    """
    Reads a text file and applies semantic chunking.
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
        chunks = semantic_chunking(text)
        return "\n\n---CHUNK---\n\n".join(chunks)
    except (FileNotFoundError, PermissionError) as e:
        raise ParserError(f"File error parsing text file {file_path}: {str(e)}")
    except Exception as e:
        raise ParserError(f"Error parsing text file {file_path}: {str(e)}")

def semantic_chunking(text: str, chunk_size: int = 700, overlap: int = 150) -> List[str]:
    """
    Splits text into chunks based on logical sections (double newlines)
    and maintains a context window (overlap). Strictly respects chunk_size.
    """
    if not text:
        return []

    sections = text.split('\n\n')
    chunks = []
    current_chunk = ""

    for section in sections:
        section = section.strip()
        if not section:
            continue

        # 1. Handle oversized sections by splitting them into chunks of chunk_size
        while len(section) > chunk_size:
            # If we have a current_chunk, save it before processing the oversized section
            if current_chunk:
                chunks.append(current_chunk)
                current_chunk = ""

            # Extract a piece of chunk_size
            piece = section[:chunk_size]
            chunks.append(piece)

            # The remaining section starts from the end of the piece minus overlap
            section = section[chunk_size - overlap:]
            # Trim section if it's still too long for the very first piece of a new chunk
            # But we are in a while loop, so it will be handled.

        # 2. Now len(section) <= chunk_size. Try to fit it into current_chunk.
        if not current_chunk:
            current_chunk = section
        elif len(current_chunk) + 2 + len(section) <= chunk_size:
            current_chunk += "\n\n" + section
        else:
            # Current chunk is full, save it and start a new one with overlap
            chunks.append(current_chunk)

            # Calculate overlap from the previous chunk
            overlap_text = current_chunk[-overlap:] if len(current_chunk) > overlap else current_chunk

            # Ensure overlap_text + "\n\n" + section fits within chunk_size
            # If it doesn't, truncate the overlap_text
            max_overlap_len = chunk_size - len(section) - 2
            if len(overlap_text) > max_overlap_len:
                overlap_text = overlap_text[max(0, len(overlap_text) - max_overlap_len):]

            if overlap_text:
                current_chunk = overlap_text + "\n\n" + section
            else:
                current_chunk = section

    if current_chunk:
        chunks.append(current_chunk)

    return chunks

def parse_pdf(file_path: str) -> str:
    """
    Reads a PDF file using pdfplumber (better for tables) and applies semantic chunking.
    """
    try:
        text = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                # Extract text
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n\n"
                
                # Extract tables and convert to text
                tables = page.extract_tables()
                for table in tables:
                    for row in table:
                        if row:
                            # Filter out None values and join
                            row_text = " | ".join(str(cell) if cell else "" for cell in row)
                            if row_text.strip():
                                text += row_text + "\n"
                    text += "\n"
        
        if not text.strip():
            # Fallback to pypdfium2 if pdfplumber fails
            import pypdfium2 as pdfium
            pdf = pdfium.PdfDocument(file_path)
            for i in range(len(pdf)):
                page = pdf.get_page(i)
                textpage = page.get_textpage()
                page_text = textpage.get_text_range()
                if page_text:
                    text += page_text + "\n\n"
        
        chunks = semantic_chunking(text)
        return "\n\n---CHUNK---\n\n".join(chunks)
    except (FileNotFoundError, PermissionError) as e:
        raise ParserError(f"File error parsing PDF file {file_path}: {str(e)}")
    except Exception as e:
        raise ParserError(f"Error parsing PDF file {file_path}: {str(e)}")

def parse_word(file_path: str) -> str:
    """
    Reads a Word file and applies semantic chunking.
    """
    try:
        doc = Document(file_path)
        text = "\n\n".join([para.text for para in doc.paragraphs])

        chunks = semantic_chunking(text)
        return "\n\n---CHUNK---\n\n".join(chunks)
    except (FileNotFoundError, PermissionError) as e:
        raise ParserError(f"File error parsing Word file {file_path}: {str(e)}")
    except Exception as e:
        # docx can raise various internal errors, wrapping them in ParserError
        raise ParserError(f"Error parsing Word file {file_path}: {str(e)}")

def parse_pptx(file_path: str) -> str:
    """
    Reads a PowerPoint file and extracts text from all slides.
    """
    try:
        prs = Presentation(file_path)
        all_text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                all_text.append(f"## Slide {slide_num}\n\n" + "\n\n".join(slide_text))
        text = "\n\n".join(all_text)
        chunks = semantic_chunking(text)
        return "\n\n---CHUNK---\n\n".join(chunks)
    except (FileNotFoundError, PermissionError) as e:
        raise ParserError(f"File error parsing PowerPoint file {file_path}: {str(e)}")
    except Exception as e:
        raise ParserError(f"Error parsing PowerPoint file {file_path}: {str(e)}")

def main():
    parser = argparse.ArgumentParser(description="Advanced Text Parsers for RAG")
    parser.add_argument("file_path", help="Path to the file to be parsed")
    args = parser.parse_args()

    file_path = args.file_path
    if not os.path.exists(file_path):
        print(f"Error: File {file_path} not found.")
        sys.exit(1)

    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext in ['.xlsx', '.xls']:
            result = parse_excel(file_path)
        elif ext == '.pdf':
            result = parse_pdf(file_path)
        elif ext == '.docx':
            result = parse_word(file_path)
        elif ext == '.pptx':
            result = parse_pptx(file_path)
        elif ext in ['.txt', '.md']:
            result = parse_text(file_path)
        else:
            print(f"Unsupported file extension: {ext}")
            sys.exit(1)
        print(result)
    except ParserError as e:
        print(f"Parsing Error: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()
